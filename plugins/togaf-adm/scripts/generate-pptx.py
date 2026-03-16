#!/usr/bin/env python3
"""
TOGAF ADM — PowerPoint Generator
Produces TOGAF-structured .pptx decks from artifact content JSON.

Usage:
    python3 generate-pptx.py --type vision --title "Architecture Vision" \
        --org "Acme Corp" --architect "Jane Smith" \
        --output "architecture-vision.pptx" --content '{"slides": [...]}'

Requirements: pip3 install python-pptx
"""

import argparse
import json
import sys
from datetime import date

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip3 install python-pptx")
    sys.exit(1)

# Brand colours
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
MED_BLUE  = RGBColor(0x2E, 0x74, 0xB5)
ORANGE    = RGBColor(0xC5, 0x5A, 0x11)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF2, 0xF7, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

DEFAULT_DECKS = {
    "vision": [
        {"layout": "title",   "title": "{title}",                   "subtitle": "{org} | {date}"},
        {"layout": "content", "title": "Executive Summary",          "bullets": ["Business context and key drivers", "Scope of this architecture engagement", "High-level target state", "Expected benefits"]},
        {"layout": "content", "title": "Architecture Scope",         "bullets": ["In scope: [define]", "Out of scope: [define]", "Assumptions and constraints"]},
        {"layout": "content", "title": "Stakeholder Summary",        "bullets": ["Key stakeholders and their concerns", "Engagement approach", "Decision-making authority"]},
        {"layout": "content", "title": "Baseline Architecture",      "bullets": ["Current state summary", "Key pain points", "Identified gaps"]},
        {"layout": "content", "title": "Target Architecture",        "bullets": ["Target state vision", "Key capabilities to develop", "Guiding principles"]},
        {"layout": "content", "title": "Gap Analysis Summary",       "bullets": ["Critical gaps", "Disposition (Retain / Replace / New)", "Priority order"]},
        {"layout": "content", "title": "Architecture Roadmap",       "bullets": ["Phase 1: Foundation (Q1–Q2)", "Phase 2: Transformation (Q3–Q4)", "Phase 3: Target State (Year 2+)"]},
        {"layout": "content", "title": "Risks and Assumptions",      "bullets": ["Top 3 risks and mitigations", "Key assumptions", "Dependencies"]},
        {"layout": "content", "title": "Next Steps",                 "bullets": ["Approve Statement of Architecture Work", "Begin Phase B — Business Architecture", "Schedule stakeholder workshops"]},
    ],
    "gap-analysis": [
        {"layout": "title",   "title": "{title}",                   "subtitle": "{org} | {date}"},
        {"layout": "content", "title": "Scope",                     "bullets": ["Architecture domain covered", "Phases analysed"]},
        {"layout": "table",   "title": "Gap Analysis",              "headers": ["ID", "Baseline", "Target", "Gap Type", "Disposition", "Priority"]},
        {"layout": "content", "title": "Key Findings",              "bullets": ["Summary of critical gaps", "Recommended dispositions", "Implementation priorities"]},
    ],
    "roadmap": [
        {"layout": "title",   "title": "{title}",                   "subtitle": "{org} | {date}"},
        {"layout": "content", "title": "Migration Strategy",        "bullets": ["Approach and rationale", "Number of transition architectures", "Sequencing logic"]},
        {"layout": "content", "title": "Architecture Roadmap",      "bullets": ["[Mermaid Gantt diagram to be embedded]", "Current → Transition 1 → Transition 2 → Target"]},
        {"layout": "content", "title": "Work Packages",             "bullets": ["Key work packages by phase", "Dependencies between packages", "Resource requirements"]},
        {"layout": "content", "title": "Next Steps",                "bullets": ["Confirm roadmap with stakeholders", "Initiate Phase G — Implementation Governance"]},
    ],
    "stakeholder": [
        {"layout": "title",   "title": "{title}",                   "subtitle": "{org} | {date}"},
        {"layout": "table",   "title": "Stakeholder Register",      "headers": ["Stakeholder", "Role", "Concerns", "Involvement", "Engagement"]},
        {"layout": "content", "title": "Power / Interest Map",      "bullets": ["High Power / High Interest: Manage closely", "High Power / Low Interest: Keep satisfied", "Low Power / High Interest: Keep informed", "Low Power / Low Interest: Monitor"]},
        {"layout": "content", "title": "Engagement Plan",           "bullets": ["Communication frequency", "Reporting format", "Escalation path"]},
    ],
}


def add_title_slide(prs, title, subtitle, org, architect, doc_date):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle.format(org=org, date=doc_date)
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)

    # Architect line
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"Prepared by: {architect}"
    run3.font.size = Pt(12)
    run3.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)

    # Orange accent bar
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(4.0), Inches(2.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    return slide


def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Light background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets[:7]):  # Max 7 bullets
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BLUE
        p.space_before = Pt(6)

    return slide


def add_table_slide(prs, title, headers, rows):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if not headers:
        return slide

    n_cols = len(headers)
    n_rows = len(rows) + 1

    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5)).table

    col_width = Inches(12.5) // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_width

    # Headers
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.size = Pt(12)

    # Rows
    for ri, row in enumerate(rows[:15]):  # Max 15 data rows
        for ci, val in enumerate(row[:n_cols]):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD6, 0xE4, 0xF0)
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].font.size = Pt(11)

    return slide


def build_deck(args, content):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    doc_date = date.today().strftime("%B %Y")
    deck_type = args.type
    slides = content.get("slides", [])

    if not slides:
        slides = DEFAULT_DECKS.get(deck_type, DEFAULT_DECKS["vision"])

    for slide_def in slides:
        layout = slide_def.get("layout", "content")
        title  = slide_def.get("title", "").format(title=args.title, org=args.org, date=doc_date)

        if layout == "title":
            subtitle = slide_def.get("subtitle", "{org} | {date}")
            add_title_slide(prs, title, subtitle, args.org, args.architect, doc_date)
        elif layout == "table":
            headers = slide_def.get("headers", content.get("table_headers", []))
            rows    = slide_def.get("rows",    content.get("table_rows", []))
            add_table_slide(prs, title, headers, rows)
        else:
            bullets = slide_def.get("bullets", ["[Content to be added]"])
            add_content_slide(prs, title, bullets)

    prs.save(args.output)
    print(f"✅ PowerPoint saved: {args.output}")


def main():
    parser = argparse.ArgumentParser(description="Generate TOGAF PowerPoint deck")
    parser.add_argument("--type",      required=True, help="Deck type (vision, gap-analysis, roadmap, stakeholder)")
    parser.add_argument("--title",     required=True, help="Deck title")
    parser.add_argument("--org",       required=True, help="Organisation name")
    parser.add_argument("--architect", required=True, help="Lead architect name")
    parser.add_argument("--output",    required=True, help="Output filename (.pptx)")
    parser.add_argument("--content",   default="{}", help="JSON content string or @file.json")
    args = parser.parse_args()

    content_str = args.content
    if content_str.startswith("@"):
        with open(content_str[1:]) as f:
            content_str = f.read()
    try:
        content = json.loads(content_str)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON — {e}")
        sys.exit(1)

    build_deck(args, content)


if __name__ == "__main__":
    main()
